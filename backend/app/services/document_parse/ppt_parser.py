"""PPT parser using python-pptx.

Extracts text slide by slide from PowerPoint files and returns structured
data with slide indices, titles, and metadata.
"""

import logging
from typing import List, Dict, Any

from pptx import Presentation

logger = logging.getLogger("course_agent.document_parse.ppt")


def parse_ppt(file_path: str) -> List[Dict[str, Any]]:
    """Extract text from a PowerPoint file slide by slide.

    Args:
        file_path: Absolute or relative path to the .ppt or .pptx file.

    Returns:
        A list of dicts, each containing:
            - slide_index (int): 1-based slide number.
            - slide_title (str): Title text from the slide (if any).
            - content (str): All text extracted from the slide.
            - metadata (dict): File-level metadata.

    Raises:
        FileNotFoundError: If the file does not exist.
    """
    results: List[Dict[str, Any]] = []

    prs = Presentation(file_path)
    total_slides = len(prs.slides)
    logger.info(f"Parsing PPT: {file_path}, {total_slides} slides")

    # Collect document-level metadata
    doc_metadata = {
        "total_slides": total_slides,
        "file_path": file_path,
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
    }

    for slide_index, slide in enumerate(prs.slides, start=1):
        try:
            # Extract title
            slide_title = ""
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_title = slide.shapes.title.text.strip()

            # Extract all text from the slide
            paragraphs = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            paragraphs.append(text)

            content = "\n".join(paragraphs)

            results.append({
                "slide_index": slide_index,
                "slide_title": slide_title,
                "content": content,
                "metadata": {
                    **doc_metadata,
                    "num_shapes": len(slide.shapes),
                },
            })

        except Exception as exc:
            logger.warning(f"Failed to extract slide {slide_index}: {exc}")
            results.append({
                "slide_index": slide_index,
                "slide_title": "",
                "content": "",
                "metadata": {
                    **doc_metadata,
                    "error": str(exc),
                },
            })

    return results
